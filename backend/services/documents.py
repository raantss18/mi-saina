"""Extraction de texte depuis des documents (PDF, Word, Excel, PowerPoint, CSV…).

Bibliothèques pures-python, légères, optionnelles : si l'une manque, on renvoie
un message d'erreur clair plutôt que de planter. Utilisé pour :
- les pièces jointes « document » (extraites avant l'envoi au modèle) ;
- la directive [READ: chemin] (lire/résumer un fichier local).
"""
import base64
import tempfile
from pathlib import Path

# Extensions lues directement comme texte brut (code, données, balisage…).
TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".conf", ".xml", ".html", ".htm", ".tex", ".bib",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash", ".c", ".h", ".cpp",
    ".hpp", ".java", ".go", ".rs", ".rb", ".php", ".sql", ".r", ".lua",
}
DOC_EXTS = {".pdf", ".docx", ".pptx", ".xlsx"}
SUPPORTED_EXTS = TEXT_EXTS | DOC_EXTS


def is_supported(name: str) -> bool:
    return Path(name).suffix.lower() in SUPPORTED_EXTS


def _pdf(p: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(p))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    return "\n".join(par.text for par in d.paragraphs)


def _pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"-- Diapositive {i} --")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for par in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in par.runs)
                    if t.strip():
                        out.append(t)
    return "\n".join(out)


def _xlsx(p: Path, max_rows: int = 200) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(p), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"-- Feuille : {ws.title} --")
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r >= max_rows:
                out.append("[… lignes suivantes tronquées]")
                break
            out.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(out)


def extract_text(path: str, max_chars: int = 20000) -> tuple[str, str | None]:
    """Renvoie (texte, erreur). `erreur` est None en cas de succès."""
    p = Path(path).expanduser()
    if not p.exists() or not p.is_file():
        return "", f"Fichier introuvable : {path}"
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            text = _pdf(p)
        elif ext == ".docx":
            text = _docx(p)
        elif ext == ".pptx":
            text = _pptx(p)
        elif ext == ".xlsx":
            text = _xlsx(p)
        elif ext in TEXT_EXTS:
            text = p.read_text(errors="replace")
        else:
            try:
                text = p.read_text()  # tentative texte brut
            except (UnicodeDecodeError, OSError):
                return "", f"Format non pris en charge : {ext or '(sans extension)'}"
    except ModuleNotFoundError as e:
        return "", (f"Lecture de {ext} indisponible : bibliothèque « {e.name} » manquante. "
                    f"Installe-la dans le venv (ex. pip install pypdf python-docx openpyxl python-pptx).")
    except Exception as e:  # noqa: BLE001
        return "", f"Lecture impossible ({ext}) : {e}"

    text = (text or "").strip()
    if not text:
        return "", "Document vide ou sans texte extractible (PDF scanné/image ?)."
    if len(text) > max_chars:
        text = text[:max_chars] + f"\n\n[… document tronqué à {max_chars} caractères]"
    return text, None


def extract_from_bytes(data: bytes, name: str, max_chars: int = 20000) -> tuple[str, str | None]:
    """Extrait le texte d'un document fourni en mémoire (pièce jointe)."""
    suffix = Path(name).suffix or ".bin"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=True) as tf:
        tf.write(data)
        tf.flush()
        return extract_text(tf.name, max_chars=max_chars)


def extract_from_base64(b64: str, name: str, max_chars: int = 20000) -> tuple[str, str | None]:
    try:
        raw = base64.b64decode(b64)
    except Exception:  # noqa: BLE001
        return "", "Pièce jointe illisible (base64 invalide)."
    return extract_from_bytes(raw, name, max_chars=max_chars)
